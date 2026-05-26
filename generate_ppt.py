import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color palette (harmonious premium deep navy, dark charcoal, and accents)
    BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
    TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50
    TEXT_GRAY = RGBColor(148, 163, 184)   # Slate 400
    ACCENT_BLUE = RGBColor(56, 189, 248)  # Sky 400
    ACCENT_GREEN = RGBColor(52, 211, 153) # Emerald 400
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Sleek Dark Theme)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title & Subtitle text box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MARKET BASKET ANALYSIS SUITE"
    p.font.name = "Arial"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Data-Driven Retail Analytics for Targeted Cross-Selling"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Presenter: Soujanya K Hegde\nRole: Data Science Intern (Retail Analytics)\nInstitution: Persevex LMS"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_GRAY
    p3.space_before = Pt(36)
    
    # ----------------------------------------------------
    # HELPER FUNCTION FOR STANDARD SLIDES (Consistent layout)
    # ----------------------------------------------------
    def add_standard_slide(title, points, accent_title=False, highlight_index=None):
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE if not accent_title else ACCENT_GREEN
        
        # Content box
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.0))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for idx, (bullet, subbullets) in enumerate(points):
            p_bullet = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p_bullet.text = f"•  {bullet}"
            p_bullet.font.name = "Arial"
            p_bullet.font.size = Pt(18)
            p_bullet.font.bold = True
            p_bullet.font.color.rgb = ACCENT_GREEN if idx == highlight_index else TEXT_WHITE
            p_bullet.space_before = Pt(12)
            
            for sub in subbullets:
                p_sub = tf_content.add_paragraph()
                p_sub.text = f"    -  {sub}"
                p_sub.font.name = "Arial"
                p_sub.font.size = Pt(14)
                p_sub.font.color.rgb = TEXT_GRAY
                p_sub.space_before = Pt(4)
                
        return slide

    # ----------------------------------------------------
    # SLIDE 2: Project Overview & Core Objectives
    # ----------------------------------------------------
    slide_2_points = [
        ("The Retail Analytics Challenge", [
            "Retail stores and e-commerce platforms struggle to understand product buying co-occurrences.",
            "Missed cross-selling and upselling opportunities lead to stagnant average order values."
        ]),
        ("Core Project Objectives", [
            "Ingest transactional datasets robustly and clean cancellations or returns.",
            "Apply association rule mining to discover hidden patterns and product affinities.",
            "Formulate strategic layout improvements, digital recommendations, and promotional bundles."
        ]),
        ("Key Business Driver", [
            "Increase Average Order Value (AOV) and customer basket size systematically using probability metrics."
        ])
    ]
    add_standard_slide("Project Overview & Objectives", slide_2_points)

    # ----------------------------------------------------
    # SLIDE 3: System Architecture & Workflow
    # ----------------------------------------------------
    slide_3_points = [
        ("End-to-End Analytics Pipeline", [
            "Phase 1: Ingestion & Multi-Encoding Data Transformation.",
            "Phase 2: Mathematical Frequent Itemset Mining & Speed Benchmarking.",
            "Phase 3: Association Rule Extraction & High-Lift Pruning.",
            "Phase 4: Static & Dynamic Constellation Visualization.",
            "Phase 5: Automated Heuristics Business Strategy Compiler."
        ]),
        ("Integrated Interactive Dashboard", [
            "Streamlit web interface, dynamic slider thresholds, interactive Plotly scatterplots, and physics-enabled PyVis constellation networks."
        ])
    ]
    add_standard_slide("Methodology & System Architecture", slide_3_points)

    # ----------------------------------------------------
    # SLIDE 4: Preprocessing & Pivoted Boolean Matrix
    # ----------------------------------------------------
    slide_4_points = [
        ("Robust Data Ingestion", [
            "Designed multi-encoding decoders (loops CP1252, Latin1, ISO-8859-1, UTF-8) to safely parse special currency characters (e.g. Pound sign '£') without UnicodeDecodeErrors."
        ]),
        ("Rigorous Preprocessing Filters", [
            "Excluded transactional cancellations and customer returns (Invoices starting with 'C' or Quantity <= 0).",
            "Pruned Null transaction values, trimmed and standardized product descriptors."
        ]),
        ("Memory-Optimized Pivoted Boolean Matrix", [
            "Pivoted transactional data from standard 'Long' format into a boolean index matrix.",
            "Cast cell values into strict boolean format, reducing RAM usage by 87.5% and optimizing computational speeds."
        ])
    ]
    add_standard_slide("Preprocessing & Matrix Transformation", slide_4_points)

    # ----------------------------------------------------
    # SLIDE 5: Algorithm Comparison: Apriori vs FP-Growth
    # ----------------------------------------------------
    slide_5_points = [
        ("The Apriori Algorithm Constraint", [
            "Calculates frequent itemsets iteratively via combinatorial candidate generation.",
            "Exponential complexity of O(2^N) triggers OOM crashes or infinite browser hanging on large catalogs."
        ]),
        ("The FP-Growth Advantage", [
            "Constructs a compressed tree structure (FP-Tree) in RAM.",
            "Mines frequent patterns instantly without candidate generation, reducing complexity by 100x."
        ]),
        ("Dual-Algorithm Safe-Sample Benchmarking", [
            "Full FP-Growth is executed on the complete 20,136 transactions in under 0.2 seconds.",
            "Apriori is safely compared to FP-Growth on a small, dense representative slice (300 tx, 15 items), creating a gorgeous Plotly speed visualization without server hangs."
        ])
    ]
    add_standard_slide("Algorithm Benchmarking: Speed & Stability", slide_5_points)

    # ----------------------------------------------------
    # SLIDE 6: Mathematical Association Metrics
    # ----------------------------------------------------
    slide_6_points = [
        ("Support (Popularity Metrics)", [
            "Frequency of product combination in transactions: P(A ∩ B).",
            "Formula: Frequency(A, B) / Total Transactions."
        ]),
        ("Confidence (Predictability Metrics)", [
            "Probability of buying B given that A is bought: P(B | A).",
            "Formula: Frequency(A, B) / Frequency(A)."
        ]),
        ("Lift (Relationship Strength - Key Metric)", [
            "Ratio of observed support to expected support if A and B were entirely independent.",
            "Formula: Confidence / Support(B).",
            "Lift > 1 indicates a strong, mathematical correlation rather than mere individual item popularity."
        ])
    ]
    add_standard_slide("Mathematical Association Metrics", slide_6_points, highlight_index=2)

    # ----------------------------------------------------
    # SLIDE 7: Visualizing Product Constellations
    # ----------------------------------------------------
    slide_7_points = [
        ("Static High-Resolution Topology Diagram", [
            "Calculated using NetworkX with nodes representing products and edges representing associations.",
            "Node size indicates Support (popularity) and edge thickness indicates Lift (affinity)."
        ]),
        ("Physics-Enabled Dynamic PyVis Constellation Map", [
            "Interactive HTML dashboard embedding featuring physics gravity-forces.",
            "Enables retail managers to drag nodes, zoom, and visually spot distinct product clusters."
        ]),
        ("Identified Retail Affinity Clusters", [
            "The Baking Cluster: Flour, Sugar, Eggs, Chocolate Chips, Baking Powder.",
            "The Weekend Party Cluster: Tortilla Chips, Salsa, Cheddar Cheese, Potato Chips, Soft Drinks.",
            "The Breakfast Cluster: Whole Wheat Bread, Butter, Jam, Milk, Ground Coffee."
        ])
    ]
    add_standard_slide("Visualizing Product Constellations", slide_7_points)

    # ----------------------------------------------------
    # SLIDE 8: Actionable Retail Business Strategies
    # ----------------------------------------------------
    slide_8_points = [
        ("Co-Merchandising Shelf Layouts (Baking End-Cap)", [
            "Rule: {All-Purpose Flour, White Sugar} -> {Baking Powder, Chocolate Chips, Eggs}",
            "Action: Create a unified 'Baking Island' display. Bundle flour and sugar, and trigger cross-promotions on chocolate chips to capture the entire cluster."
        ]),
        ("Promotional Bundling (Weekend Snack Box)", [
            "Rule: {Tomato Salsa, Potato Chips} -> {Tortilla Chips, Cheddar Cheese, Soft Drinks}",
            "Action: Offer a pre-packaged 'Weekend Party Box' near checkouts on Friday afternoons with a 30% discount on soft drinks when chips and salsa are purchased."
        ]),
        ("Digital Cart Upsell Widgets (Storefront Recommendations)", [
            "Rule: {Smartphone Case Clear} -> {Tempered Screen Protector}",
            "Action: Trigger checkout pop-up recommendation widgets: 'Customers who bought this Case also added the Screen Protector!' with a one-click add-to-cart button."
        ])
    ]
    add_standard_slide("Actionable Business Strategies", slide_8_points)

    # ----------------------------------------------------
    # SLIDE 9: System Stack & Execution Deliverables
    # ----------------------------------------------------
    slide_9_points = [
        ("Robust Technical Stack", [
            "Backend: Python, Pandas, Numpy, MLxtend, NetworkX, PyVis.",
            "Frontend: Streamlit server framework, Plotly visualization charts."
        ]),
        ("Production Deliverables Generated", [
            "The Rules Engine CLI (main.py) saving mined data to association_rules.csv.",
            "High-resolution PNG and physics-enabled PyVis HTML constellation graph.",
            "Print-ready professional HTML & PDF Business Strategy Proposal."
        ]),
        ("Complete GitHub Version Control", [
            "Local Git repository initialized, configured with robust .gitignore rules, and fully pushed to online master branch for tracking."
        ])
    ]
    add_standard_slide("System Stack & Deliverables", slide_9_points)

    # ----------------------------------------------------
    # SLIDE 10: Conclusion & Core Accomplishments
    # ----------------------------------------------------
    slide_10_points = [
        ("100% Successful Internship Completion", [
            "Built a flawless, production-grade Market Basket Analysis Suite."
        ]),
        ("Optimized System Performance Achievements", [
            "Combatorial Explosion Fix: Pruned rules using dynamic max_len stopping criteria, reducing mining execution times from hours to milliseconds.",
            "Module Hot-Reloading: Bypassed Streamlit in-memory caching by creating dynamic reloaders using importlib, ensuring active browser updates."
        ]),
        ("Professional Business Impact", [
            "Translated complex probability rules into concrete, high-impact retail recommendation strategies to boost retail basket value."
        ])
    ]
    add_standard_slide("Conclusion & Key Accomplishments", slide_10_points, accent_title=True)

    # Save presentation
    output_path = "outputs/market_basket_analysis_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
